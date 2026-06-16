#!/usr/bin/env python3
"""Content-agnostic python-pptx renderer for the pptx MCP server.

Reads a single JSON job from stdin and prints a single JSON result to stdout.
Never prints anything else to stdout (the Node parent parses stdout as JSON);
diagnostics go to stderr.

Jobs:
  {"action": "render", "spec": {...}, "templatePath": str|None, "outPath": str}
      -> {"path": str, "slides": int}  on success
  {"action": "inspect", "templatePath": str}
      -> {"layouts": [{"index": int, "name": str}]}
  {"action": "merge", "inputs": [str, ...], "outPath": str}
      -> {"path": str, "slides": int}  appends slides of inputs[1:] onto inputs[0]
  any failure -> {"error": str}

The renderer maps a small, stable layout vocabulary onto whatever layouts the
chosen base presentation provides (default template, or an uploaded one). It
degrades gracefully: if a named layout is missing it falls back to a sensible
index, so an arbitrary uploaded template still works.
"""

import json
import sys
import os


def _emit(obj):
    sys.stdout.write(json.dumps(obj))
    sys.stdout.flush()


def _fail(msg):
    _emit({"error": str(msg)})
    sys.exit(0)


def _load_pptx():
    try:
        from pptx import Presentation  # noqa: F401
        from pptx.util import Pt, Emu  # noqa: F401
        from pptx.dml.color import RGBColor  # noqa: F401
        return sys.modules["pptx"]
    except Exception as e:  # pragma: no cover - import guard
        _fail(f"python-pptx not available: {e}")


# Map our stable layout names to typical python-pptx default-template indices.
# These are best-effort; _pick_layout falls back if an index is out of range.
_LAYOUT_HINTS = {
    "title": 0,          # Title Slide
    "title_content": 1,  # Title and Content
    "section": 2,        # Section Header
    "blank": 6,          # Blank
}


def _pick_layout(prs, name):
    layouts = prs.slide_layouts
    n = len(layouts)
    if n == 0:
        raise ValueError("template has no slide layouts")
    # 1) try to match by layout name (case-insensitive contains)
    wanted = (name or "title_content").lower()
    name_keywords = {
        "title": ["title slide", "title"],
        "title_content": ["title and content", "content"],
        "section": ["section"],
        "blank": ["blank"],
    }.get(wanted, [])
    for kw in name_keywords:
        for layout in layouts:
            try:
                if kw in (layout.name or "").lower():
                    return layout
            except Exception:
                pass
    # 2) fall back to the hinted index, clamped to range
    idx = _LAYOUT_HINTS.get(wanted, 1)
    if idx >= n:
        idx = min(1, n - 1)
    return layouts[idx]


def _set_font(run, font_name=None, size_pt=None, color_hex=None):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    if font_name:
        run.font.name = font_name
    if size_pt:
        run.font.size = Pt(size_pt)
    if color_hex:
        try:
            run.font.color.rgb = RGBColor.from_string(color_hex.lstrip("#"))
        except Exception:
            pass


def _norm_bullets(bullets):
    out = []
    for b in bullets or []:
        if isinstance(b, str):
            out.append({"text": b, "level": 0})
        elif isinstance(b, dict) and "text" in b:
            out.append({"text": str(b["text"]), "level": int(b.get("level", 0) or 0)})
    return out


def _add_title(slide, text, theme):
    if not text:
        return
    title_ph = None
    try:
        title_ph = slide.shapes.title
    except Exception:
        title_ph = None
    if title_ph is not None:
        title_ph.text = text
        try:
            run = title_ph.text_frame.paragraphs[0].runs[0]
            _set_font(run, theme.get("titleFont"), None, theme.get("accentColor"))
        except Exception:
            pass


def _body_placeholder(slide):
    # Prefer an explicit body/content placeholder; else first non-title placeholder.
    from pptx.util import Pt  # noqa: F401
    title_id = None
    try:
        if slide.shapes.title is not None:
            title_id = slide.shapes.title.placeholder_format.idx
    except Exception:
        title_id = None
    candidate = None
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == title_id:
                continue
        except Exception:
            pass
        candidate = ph
        break
    return candidate


def _add_bullets(slide, bullets, theme):
    bullets = _norm_bullets(bullets)
    if not bullets:
        return
    body = _body_placeholder(slide)
    if body is None or not body.has_text_frame:
        return
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b["text"]
        p.level = max(0, min(4, b["level"]))
        try:
            _set_font(p.runs[0], theme.get("bodyFont"), None, None)
        except Exception:
            pass


def _add_subtitle(slide, text):
    if not text:
        return
    for ph in slide.placeholders:
        try:
            name = (ph.name or "").lower()
        except Exception:
            name = ""
        if "subtitle" in name and ph.has_text_frame:
            ph.text = text
            return
    # fallback: use the body placeholder if no explicit subtitle
    body = _body_placeholder(slide)
    if body is not None and body.has_text_frame:
        body.text = text


def _add_image(slide, image_path):
    from pptx.util import Inches
    if not image_path or not os.path.exists(image_path):
        return
    try:
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.8), height=Inches(4))
    except Exception:
        pass


def _add_notes(slide, notes):
    if not notes:
        return
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except Exception:
        pass


def _apply_background(slide, theme):
    color = theme.get("backgroundColor")
    if not color:
        return
    from pptx.dml.color import RGBColor
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
    except Exception:
        pass


# ----------------------------------------------------------------------------
# Self-drawn layout engine
#
# The placeholder-based helpers above keep working for legacy specs, but they
# inherit the template's centered/top-anchored look. The functions below draw
# everything with absolute-positioned shapes so we control alignment, accent
# blocks, columns, full-bleed images and big-number emphasis — i.e. real
# layout variety instead of one templated skeleton.
#
# Canvas assumes a 16:9 deck: 13.333in x 7.5in.
# ----------------------------------------------------------------------------

EMU_PER_IN = 914400
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def _hex(color_hex, default="FFFFFF"):
    from pptx.dml.color import RGBColor
    try:
        return RGBColor.from_string((color_hex or default).lstrip("#"))
    except Exception:
        return RGBColor.from_string(default)


def _is_dark(color_hex):
    """Rough luminance test so text contrasts with the background."""
    try:
        h = (color_hex or "FFFFFF").lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) < 140
    except Exception:
        return False


def _body_color(theme):
    """Readable body text color for the current background."""
    return "ECECEC" if _is_dark(theme.get("backgroundColor")) else "1A1A1A"


def _rect(slide, x, y, w, h, fill_hex, line=False):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex(fill_hex)
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _textbox(slide, x, y, w, h, text, *, size, color, font=None, bold=False,
             align="left", anchor="top", line_spacing=None):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                              "bottom": MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.TOP)
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    if line_spacing:
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _hex(color)
    if font:
        run.font.name = font
    return tb


def _accent_heading(slide, theme, text, *, x=0.9, y=0.7, w=11.5):
    """Left-aligned heading with an accent bar to its left — replaces the
    centered placeholder title for content/section slides."""
    accent = theme.get("accentColor") or "2E74B5"
    # accent bar
    _rect(slide, x, y + 0.05, 0.14, 0.95, accent)
    _textbox(slide, x + 0.32, y, w, 1.1, text,
             size=32, color=accent, font=theme.get("titleFont"), bold=True,
             align="left", anchor="middle")


def _draw_bullets(slide, bullets, theme, *, x=0.95, y=2.0, w=11.4, h=4.8, size=20):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    bullets = _norm_bullets(bullets)
    if not bullets:
        return
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    body_hex = _body_color(theme)
    accent = theme.get("accentColor") or "2E74B5"
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        lvl = max(0, min(4, b["level"]))
        p.level = lvl
        try:
            p.space_after = Pt(10)
            p.line_spacing = 1.15
        except Exception:
            pass
        # marker run (accent dot/dash) + text run
        marker = p.add_run()
        marker.text = ("›  " if lvl == 0 else "–  ")
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.color.rgb = _hex(accent)
        if theme.get("bodyFont"):
            marker.font.name = theme.get("bodyFont")
        run = p.add_run()
        run.text = b["text"]
        run.font.size = Pt(size - lvl * 2)
        run.font.color.rgb = _hex(body_hex)
        if theme.get("bodyFont"):
            run.font.name = theme.get("bodyFont")


def _full_bleed_image(slide, image_path):
    from pptx.util import Inches
    if not image_path or not os.path.exists(image_path):
        return False
    try:
        slide.shapes.add_picture(image_path, Inches(0), Inches(0),
                                 width=Inches(SLIDE_W_IN), height=Inches(SLIDE_H_IN))
        return True
    except Exception:
        return False


def _scrim(slide, hex_color="000000", *, y=0.0, h=SLIDE_H_IN, alpha_dark=True):
    """A semi-transparent-ish band so text on an image stays readable.
    python-pptx can't do true alpha easily; use a solid band low-opacity-like
    by drawing a dark rectangle (works well for bottom/side scrims)."""
    _rect(slide, 0, y, SLIDE_W_IN, h, hex_color)


def _image_box(slide, image_path, x, y, w, h):
    from pptx.util import Inches
    if not image_path or not os.path.exists(image_path):
        return False
    try:
        slide.shapes.add_picture(image_path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        return True
    except Exception:
        return False


# Layouts that the self-drawn engine handles. Anything else falls back to the
# legacy placeholder path.
_SELF_DRAWN = {
    "title", "section", "content", "title_content",
    "two_column", "left_text_right_image", "right_text_left_image",
    "big_number", "full_bleed_image", "quote",
}


def _draw_cover(slide, theme, title, subtitle):
    accent = theme.get("accentColor") or "2E74B5"
    # big left-aligned title block, accent underline, subtitle below
    _textbox(slide, 0.95, 2.4, 11.4, 2.0, title or "",
             size=54, color=accent, font=theme.get("titleFont"), bold=True,
             align="left", anchor="bottom")
    _rect(slide, 1.0, 4.45, 3.6, 0.12, accent)
    if subtitle:
        _textbox(slide, 0.98, 4.7, 11.0, 1.2, subtitle,
                 size=22, color=_body_color(theme), font=theme.get("bodyFont"),
                 align="left", anchor="top")


def _draw_section(slide, theme, title, subtitle, index=None):
    accent = theme.get("accentColor") or "2E74B5"
    if index is not None:
        _textbox(slide, 0.95, 1.6, 4.0, 1.6, str(index).zfill(2),
                 size=96, color=accent, font=theme.get("titleFont"), bold=True,
                 align="left", anchor="top")
    _rect(slide, 1.0, 3.7, 2.4, 0.14, accent)
    _textbox(slide, 0.95, 3.95, 11.4, 1.6, title or "",
             size=40, color=_body_color(theme), font=theme.get("titleFont"), bold=True,
             align="left", anchor="top")
    if subtitle:
        _textbox(slide, 0.98, 5.4, 11.0, 1.0, subtitle,
                 size=20, color=_body_color(theme), font=theme.get("bodyFont"),
                 align="left", anchor="top")


def _draw_big_number(slide, theme, title, number, caption):
    accent = theme.get("accentColor") or "2E74B5"
    if title:
        _accent_heading(slide, theme, title)
    _textbox(slide, 0.9, 2.3, 11.5, 2.6, number or "",
             size=130, color=accent, font=theme.get("titleFont"), bold=True,
             align="left", anchor="middle")
    if caption:
        _textbox(slide, 0.98, 5.2, 11.0, 1.4, caption,
                 size=24, color=_body_color(theme), font=theme.get("bodyFont"),
                 align="left", anchor="top")


def _draw_quote(slide, theme, text, attribution):
    accent = theme.get("accentColor") or "2E74B5"
    _textbox(slide, 1.2, 1.4, 1.6, 1.6, "\u201C",
             size=120, color=accent, font=theme.get("titleFont"), bold=True, align="left")
    _textbox(slide, 1.3, 2.6, 10.7, 3.0, text or "",
             size=34, color=_body_color(theme), font=theme.get("titleFont"), bold=False,
             align="left", anchor="top", line_spacing=1.2)
    if attribution:
        _textbox(slide, 1.32, 5.8, 10.0, 0.8, "— " + attribution,
                 size=20, color=accent, font=theme.get("bodyFont"), align="left")


def _draw_two_column(slide, theme, title, left, right):
    if title:
        _accent_heading(slide, theme, title)
    _draw_bullets(slide, left, theme, x=0.95, y=2.0, w=5.5, h=4.8)
    _draw_bullets(slide, right, theme, x=6.9, y=2.0, w=5.5, h=4.8)


def _draw_text_image(slide, theme, title, bullets, image_path, *, image_right=True):
    if title:
        _accent_heading(slide, theme, title)
    if image_right:
        _draw_bullets(slide, bullets, theme, x=0.95, y=2.0, w=6.0, h=4.8)
        _image_box(slide, image_path, 7.3, 2.0, 5.2, 4.4)
    else:
        _image_box(slide, image_path, 0.8, 2.0, 5.2, 4.4)
        _draw_bullets(slide, bullets, theme, x=6.5, y=2.0, w=6.0, h=4.8)


def _draw_content(slide, theme, title, bullets, image_path):
    if title:
        _accent_heading(slide, theme, title)
    if image_path and os.path.exists(image_path) and not bullets:
        _image_box(slide, image_path, 2.4, 2.0, 8.5, 4.6)
    else:
        _draw_bullets(slide, bullets, theme, x=0.95, y=2.0, w=11.4, h=4.8)
        if image_path and os.path.exists(image_path):
            # small supporting image lower-right
            _image_box(slide, image_path, 8.6, 4.3, 3.8, 2.6)


def _render_self_drawn(slide, theme, s, *, section_index=None):
    """Render one slide using the self-drawn engine. Returns True if handled."""
    layout = (s.get("layout") or "content").lower()
    title = s.get("title")
    if layout == "title":
        _draw_cover(slide, theme, title, s.get("subtitle"))
    elif layout == "section":
        _draw_section(slide, theme, title, s.get("subtitle"), index=s.get("index", section_index))
    elif layout == "big_number":
        _draw_big_number(slide, theme, title, s.get("number") or s.get("subtitle"), s.get("caption"))
    elif layout == "quote":
        _draw_quote(slide, theme, s.get("text") or title, s.get("attribution"))
    elif layout == "two_column":
        _draw_two_column(slide, theme, title, s.get("left") or s.get("bullets"), s.get("right"))
    elif layout in ("left_text_right_image",):
        _draw_text_image(slide, theme, title, s.get("bullets"), s.get("imagePath"), image_right=True)
    elif layout in ("right_text_left_image",):
        _draw_text_image(slide, theme, title, s.get("bullets"), s.get("imagePath"), image_right=False)
    elif layout == "full_bleed_image":
        ok_img = _full_bleed_image(slide, s.get("imagePath"))
        if title:
            # bottom scrim band + title for readability
            _rect(slide, 0, SLIDE_H_IN - 1.9, SLIDE_W_IN, 1.9, "0B0B0B")
            _textbox(slide, 0.9, SLIDE_H_IN - 1.7, 11.5, 1.4, title,
                     size=40, color="FFFFFF", font=theme.get("titleFont"), bold=True,
                     align="left", anchor="middle")
    else:  # "content" / "title_content"
        _draw_content(slide, theme, title, s.get("bullets"), s.get("imagePath"))
    _add_notes(slide, s.get("notes"))
    return True


def action_render(job):
    pptx = _load_pptx()
    from pptx import Presentation

    spec = job.get("spec") or {}
    template_path = job.get("templatePath")
    out_path = job.get("outPath")
    if not out_path:
        raise ValueError("outPath is required")

    if template_path:
        if not os.path.exists(template_path):
            raise ValueError(f"templatePath does not exist: {template_path}")
        prs = Presentation(template_path)
        # Remove any slides that ship inside the template so we start clean,
        # while keeping its masters/layouts/theme.
        xml_slides = prs.slides._sldIdLst
        for sldId in list(xml_slides):
            xml_slides.remove(sldId)
    else:
        prs = Presentation()

    theme = (spec.get("theme") or {}) if not template_path else {}

    # Ensure a 16:9 canvas when building from scratch (matches the layout engine).
    if not template_path:
        try:
            from pptx.util import Inches
            prs.slide_width = Inches(SLIDE_W_IN)
            prs.slide_height = Inches(SLIDE_H_IN)
        except Exception:
            pass

    slides_spec = spec.get("slides") or []
    title = spec.get("title")
    subtitle = spec.get("subtitle")
    want_title_slide = spec.get("titleSlide", True)
    count = 0
    section_counter = 0

    blank = _pick_layout(prs, "blank")

    def _new_slide():
        sl = prs.slides.add_slide(blank)
        _apply_background(sl, theme)
        return sl

    # Optional dedicated cover slide.
    if want_title_slide and title and (not slides_spec or slides_spec[0].get("layout") != "title"):
        slide = _new_slide()
        _render_self_drawn(slide, theme, {"layout": "title", "title": title, "subtitle": subtitle})
        count += 1

    for s in slides_spec:
        layout = (s.get("layout") or "content").lower()
        slide = _new_slide()
        if layout == "section":
            section_counter += 1
        if layout in _SELF_DRAWN:
            _render_self_drawn(slide, theme, s, section_index=section_counter)
        else:
            # legacy placeholder fallback for unknown layout names
            from_layout = _pick_layout(prs, s.get("layout") or "title_content")
            # replace the blank slide with a placeholder-based one
            slide = prs.slides.add_slide(from_layout)
            _apply_background(slide, theme)
            _add_title(slide, s.get("title"), theme)
            if s.get("subtitle"):
                _add_subtitle(slide, s.get("subtitle"))
            _add_bullets(slide, s.get("bullets"), theme)
            _add_image(slide, s.get("imagePath"))
            _add_notes(slide, s.get("notes"))
        count += 1

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return {"path": out_path, "slides": count}


def action_inspect(job):
    _load_pptx()
    from pptx import Presentation
    template_path = job.get("templatePath")
    if not template_path or not os.path.exists(template_path):
        raise ValueError(f"templatePath does not exist: {template_path}")
    prs = Presentation(template_path)
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        try:
            layouts.append({"index": i, "name": layout.name})
        except Exception:
            layouts.append({"index": i, "name": f"layout {i}"})
    return {"layouts": layouts}


def _copy_slide_into(dest_prs, src_slide):
    """Deep-copy one slide (XML + its part relationships) from a source
    presentation into dest_prs. python-pptx has no public 'copy slide' API,
    so we clone the slide part's element and re-attach external relationships
    (images, media) by copying the referenced parts and rewriting r:embed/r:id.

    This avoids the duplicate-slideN.xml / dropped-slide corruption that occurs
    when naively reusing a finished deck as a template.
    """
    import copy
    from pptx.oxml.ns import qn
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    # Use the first available layout in the destination as the new slide's layout.
    dest_layout = dest_prs.slide_layouts[0]
    new_slide = dest_prs.slides.add_slide(dest_layout)

    # Remove every shape the blank layout seeded, so we start from an empty slide
    # and bring over exactly the source slide's shape tree.
    spTree = new_slide.shapes._spTree
    for shape in list(new_slide.shapes):
        spTree.remove(shape._element)

    # Deep-copy each top-level shape element from the source slide.
    src_part = src_slide.part
    dest_part = new_slide.part
    # Map old rId -> new rId for any relationships referenced by the copied shapes.
    rId_map = {}

    def _clone_related_part(old_rId):
        if old_rId in rId_map:
            return rId_map[old_rId]
        rel = src_part.rels[old_rId]
        if rel.is_external:
            new_rId = dest_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            src_target = rel.target_part
            # copy the binary part (e.g. image) into the dest package, letting
            # next_partname fill the numbering into a '%d' template. Preserve the
            # source extension so content-type detection stays correct.
            src_name = str(src_target.partname)
            ext = src_name.rsplit(".", 1)[-1] if "." in src_name else "bin"
            partname = dest_part.package.next_partname("/ppt/media/image%d." + ext)
            new_part = Part(partname, src_target.content_type, dest_part.package, src_target.blob)
            new_rId = dest_part.relate_to(new_part, rel.reltype)
        rId_map[old_rId] = new_rId
        return new_rId

    for shp in src_slide.shapes._spTree.iterchildren():
        tag = shp.tag
        # skip the non-shape children that the spTree starts with
        if tag.endswith("}nvGrpSpPr") or tag.endswith("}grpSpPr"):
            continue
        cloned = copy.deepcopy(shp)
        # rewrite any r:embed / r:link / r:id attributes to remapped rIds
        for el in cloned.iter():
            for attr in (qn("r:embed"), qn("r:link"), qn("r:id")):
                val = el.get(attr)
                if val and val in src_part.rels:
                    el.set(attr, _clone_related_part(val))
        spTree.append(cloned)

    # Bring over speaker notes if present.
    try:
        if src_slide.has_notes_slide and src_slide.notes_slide.notes_text_frame.text:
            new_slide.notes_slide.notes_text_frame.text = (
                src_slide.notes_slide.notes_text_frame.text
            )
    except Exception:
        pass


def action_merge(job):
    """Merge several .pptx files (in order) into one deck.

    Job: {"action":"merge", "inputs":[path, ...], "outPath": path}
    The first input seeds the base (preserving its theme/size); subsequent
    inputs' slides are appended in order. Returns {"path", "slides"}.
    """
    _load_pptx()
    from pptx import Presentation

    inputs = job.get("inputs") or []
    out_path = job.get("outPath")
    if not isinstance(inputs, list) or len(inputs) < 1:
        raise ValueError("inputs must be a non-empty list of .pptx paths")
    if not out_path:
        raise ValueError("outPath is required")
    for p in inputs:
        if not p or not os.path.exists(p):
            raise ValueError(f"input does not exist: {p}")

    # Seed from the first deck so slide size + theme carry over.
    base = Presentation(inputs[0])
    for extra_path in inputs[1:]:
        src = Presentation(extra_path)
        for slide in src.slides:
            _copy_slide_into(base, slide)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    base.save(out_path)
    total = len(base.slides._sldIdLst)
    return {"path": out_path, "slides": total}


def main():
    raw = sys.stdin.read()
    try:
        job = json.loads(raw)
    except Exception as e:
        _fail(f"invalid job JSON: {e}")
        return
    action = job.get("action")
    try:
        if action == "render":
            _emit(action_render(job))
        elif action == "inspect":
            _emit(action_inspect(job))
        elif action == "merge":
            _emit(action_merge(job))
        else:
            _fail(f"unknown action: {action}")
    except Exception as e:
        _fail(e)


if __name__ == "__main__":
    main()
